# -*- coding: utf-8 -*-
"""
Générateur PPTX Barbier Immobilier — Dossier de présentation v2
Chevrons en freeform, logo réel, fidèle à la trame
"""
import io, base64, requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Couleurs ──────────────────────────────────────────────────
TEAL      = RGBColor(0x16, 0x70, 0x8B)
TEAL_DARK = RGBColor(0x0D, 0x55, 0x70)
ORANGE    = RGBColor(0xF0, 0x79, 0x5B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1B, 0x3A, 0x5C)
GRAY_LIGHT= RGBColor(0xF3, 0xF4, 0xF6)
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)

In = Inches
SW = In(13.33)
SH = In(7.5)

def _prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _b64_buf(b64):
    return io.BytesIO(base64.b64decode(b64))

def _add_rect(slide, l, t, w, h, rgb=None):
    sp = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    sp.line.fill.background()
    if rgb:
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb
    else:
        sp.fill.background()
    return sp

def _add_txt(slide, text, l, t, w, h, size=14, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, italic=False, name="Calibri"):
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = name
    r.font.italic = italic
    return tb

def _add_img(slide, src, l, t, w, h):
    if src is None: return None
    if isinstance(src, str):
        buf = _b64_buf(src)
    elif isinstance(src, (bytes, bytearray)):
        buf = io.BytesIO(src)
    else:
        src.seek(0); buf = src
    try:
        return slide.shapes.add_picture(buf, int(l), int(t), int(w), int(h))
    except Exception:
        return None

def _add_chevron_shape(slide, l, t, w, h, rgb):
    """Chevron > via freeform (4 points)"""
    from pptx.util import Emu
    # Points du chevron > :  haut-gauche, milieu-droit, bas-gauche, milieu (encoche)
    # On utilise add_shape avec MSO_SHAPE_TYPE = chevron (13)
    # MAIS on remplace l'XML pour forcer le bon preset
    sp = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    # Modifier le preset en "chevron" via XML
    sp_elem = sp._element
    spPr = sp_elem.find(qn('p:spPr'))
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'rightArrow')
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb
    sp.line.fill.background()
    return sp

def _get_photo(val):
    if not val: return None
    try:
        if val.startswith('data:'):
            _, b64 = val.split(',', 1)
            return io.BytesIO(base64.b64decode(b64))
        else:
            r = requests.get(val, timeout=8, headers={'User-Agent':'BarbierImmo/1.0'})
            r.raise_for_status()
            return io.BytesIO(r.content)
    except Exception:
        return None

def _pfmt(v):
    try: return f"{int(float(v)):,} €".replace(',', '\u202f')
    except: return str(v)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — COUVERTURE
# ══════════════════════════════════════════════════════════════
def slide_couverture(prs, d, A):
    s = _blank(prs)
    _add_rect(s, 0, 0, SW, SH, TEAL)  # fond teal

    # Photos dans les zones chevrons (zone gauche-centre)
    photo = _get_photo(d.get('Photo bien (URL)','') or d.get('photo_url',''))
    if photo:
        _add_img(s, photo, In(2.3), In(0.1), In(3.5), SH*0.95)
        photo.seek(0)
        _add_img(s, photo, In(5.2), In(1.3), In(2.8), SH*0.72)
    elif A.get('COUVERTURE_BG_B64'):
        _add_img(s, A['COUVERTURE_BG_B64'], In(2.3), In(0.1), In(5.8), SH*0.95)

    # Logo + tagline haut gauche
    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(0.25), In(0.15), In(1.4), In(0.68))
    _add_txt(s, "Votre projet devient le nôtre",
             In(0.25), In(0.9), In(2.6), In(0.35), size=9, color=WHITE)

    # Textes droite
    type_b  = d.get('Type de bien','') or '[Type de bien]'
    adresse = d.get('Adresse','') or ''
    ville   = d.get('Ville','') or ''
    adr_full = (adresse + ", " + ville).strip(', ')
    _add_txt(s, "Dossier de présentation",
             In(8.3), In(1.2), In(4.8), In(1.1),
             size=28, color=WHITE, name="Calibri Light")
    _add_txt(s, type_b,
             In(8.3), In(2.4), In(4.8), In(0.8), size=20, color=WHITE)
    _add_txt(s, "[" + adr_full + "]",
             In(9.5), In(6.3), In(3.8), In(0.6), size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ══════════════════════════════════════════════════════════════
def slide_sommaire(prs, d, A):
    s = _blank(prs)
    _add_rect(s, 0, 0, SW, SH, WHITE)
    _add_rect(s, 0, 0, In(0.15), SH, TEAL)  # bande teal gauche

    # Chevrons titre
    _add_chevron_shape(s, In(1.5), In(0.2), In(0.5), In(0.5), ORANGE)
    _add_txt(s, "SOMMAIRE", In(2.2), In(0.18), In(9.0), In(0.7),
             size=28, bold=True, color=DARK)
    _add_chevron_shape(s, In(10.5), In(0.2), In(0.5), In(0.5), TEAL_DARK)

    items = [
        "Présentation de l'agence",
        "Secteur géographique",
        "Présentation du bien",
        "  • Description du bien",
        "  • Informations cadastrales et plan",
        "Valeur locative",
        "Conclusion",
    ]
    for i, item in enumerate(items):
        y = In(1.4) + i * In(0.65)
        sub = item.startswith('  •')
        clr = ORANGE if i < 2 else TEAL if i < 5 else DARK
        if not sub:
            _add_chevron_shape(s, In(1.4), y + In(0.1), In(0.35), In(0.35), clr)
        _add_txt(s, item.strip(), In(2.0), y, In(9.5), In(0.55),
                 size=15, color=DARK, bold=(not sub))

    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(0.3), In(7.0), In(1.3), In(0.42))

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — PRÉSENTATION AGENCE
# ══════════════════════════════════════════════════════════════
def slide_agence(prs, d, A):
    s = _blank(prs)
    _add_rect(s, 0, 0, SW, SH, WHITE)
    _add_rect(s, 0, 0, In(4.0), In(3.6), TEAL)

    _add_txt(s,
        "Depuis 33 ans, barbier immobilier vous accompagne pour la "
        "réalisation de votre projet immobilier professionnel, la reprise "
        "ou la cession de votre entreprise dans le Morbihan, Vannes.",
        In(0.3), In(0.3), In(3.5), In(3.0),
        size=13, color=WHITE)

    if A.get('VOILIER_B64'):
        _add_img(s, A['VOILIER_B64'], In(0), In(3.6), In(4.0), In(3.9))

    # Logo centre
    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(4.6), In(2.3), In(3.5), In(1.6))
    _add_txt(s, "Votre projet devient le nôtre",
             In(4.4), In(4.1), In(4.2), In(0.5),
             size=13, color=DARK, name="Calibri Light")

    # Droite : photo Vannes + métiers
    if A.get('VANNES_VIEUX_B64'):
        _add_img(s, A['VANNES_VIEUX_B64'], In(8.8), In(0), In(4.5), In(3.6))
    _add_rect(s, In(8.8), In(3.6), In(4.5), In(3.9), TEAL)

    metiers = ["Estimation", "Vente", "Location", "Cession d'entreprise"]
    for i, m in enumerate(metiers):
        y = In(4.1) + i * In(0.72)
        indent = i * In(0.25)
        _add_chevron_shape(s, In(9.1)+indent, y+In(0.1), In(0.3), In(0.3), ORANGE)
        _add_txt(s, m, In(9.55)+indent, y, In(3.6), In(0.55),
                 size=14, bold=True, color=WHITE)

    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(0.2), In(7.08), In(1.2), In(0.38))

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — ÉQUIPE + EXPERTISE
# ══════════════════════════════════════════════════════════════
def slide_equipe(prs, d, A):
    s = _blank(prs)
    _add_rect(s, 0, 0, SW, SH, WHITE)

    if A.get('EQUIPE_B64'):
        _add_img(s, A['EQUIPE_B64'], In(0), In(0), In(4.5), In(3.5))
    if A.get('CHANTIER_B64'):
        _add_img(s, A['CHANTIER_B64'], In(8.9), In(0), In(4.4), In(3.5))

    _add_chevron_shape(s, In(4.7), In(0.25), In(0.4), In(0.4), ORANGE)
    _add_txt(s, "Notre équipe", In(5.3), In(0.2), In(4.0), In(0.6),
             size=22, color=TEAL)
    _add_chevron_shape(s, In(8.7), In(0.25), In(0.4), In(0.4), TEAL_DARK)

    _add_txt(s,
        "Une équipe professionnelle, compétente, bienveillante et engagée pour vous accompagner!\n\n"
        "Communication, réactivité et transparence, nous travaillerons ensemble pour mener à bien votre projet!",
        In(4.7), In(1.1), In(4.0), In(2.2), size=11, color=DARK)

    # Bas gauche : expertise
    _add_chevron_shape(s, In(0.3), In(3.7), In(0.4), In(0.4), ORANGE)
    _add_txt(s, "Notre expertise", In(0.85), In(3.65), In(4.0), In(0.6),
             size=20, color=TEAL)
    _add_chevron_shape(s, In(5.0), In(3.7), In(0.4), In(0.4), TEAL_DARK)

    expertise = [
        "Connaissance du marché",
        "Appréciation juridique et administrative,",
        "Étude personnalisée,  Prospection,",
        "Portefeuille qualifié,  Accompagnement,",
        "Rédaction des documents contractuels,",
        "Mise en relation,",
    ]
    for i, item in enumerate(expertise):
        y = In(4.4) + i * In(0.35)
        _add_chevron_shape(s, In(0.35), y+In(0.07), In(0.22), In(0.22), ORANGE)
        _add_txt(s, item, In(0.7), y, In(4.5), In(0.38), size=10, color=DARK)

    _add_txt(s,
        "Mais surtout une confiance réciproque nous permettant d'avancer sereinement!",
        In(0.3), In(6.9), In(5.2), In(0.5),
        size=11, color=ORANGE, bold=False)

    if A.get('BOUTIQUE_B64'):
        _add_img(s, A['BOUTIQUE_B64'], In(5.5), In(3.5), In(3.3), In(3.95))

    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(9.5), In(4.3), In(2.6), In(1.2))
    _add_txt(s,
        "Agence historiquement implantée à VANNES,\nnous vous accompagnons sur tout le MORBIHAN!",
        In(9.2), In(5.7), In(4.0), In(1.5), size=11, color=DARK)

    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(0.2), In(7.08), In(1.2), In(0.38))

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — SECTEUR GÉOGRAPHIQUE
# ══════════════════════════════════════════════════════════════
def slide_secteur(prs, d, A, map_buf=None):
    s = _blank(prs)
    _add_rect(s, 0, 0, SW, SH, WHITE)
    ville = d.get('Ville','') or 'Vannes'

    _add_chevron_shape(s, In(1.3), In(0.12), In(0.5), In(0.5), ORANGE)
    _add_txt(s, "Secteur géographique", In(2.0), In(0.1), In(9.5), In(0.65),
             size=26, bold=True, color=DARK)
    _add_chevron_shape(s, In(11.0), In(0.12), In(0.5), In(0.5), TEAL_DARK)

    _add_rect(s, 0, In(0.82), SW, In(0.95), TEAL)
    _add_txt(s, "[" + ville + "]", In(0.5), In(0.87), In(5.0), In(0.75),
             size=17, color=WHITE)

    # Col 1 : Bretagne
    if A.get('BRETAGNE_B64'):
        _add_img(s, A['BRETAGNE_B64'], In(0.3), In(2.0), In(3.8), In(3.8))
    _add_txt(s, "MORBIHAN", In(0.3), In(6.05), In(3.8), In(0.5),
             size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Col 2 : carte
    if map_buf:
        map_buf.seek(0)
        _add_img(s, map_buf, In(4.3), In(2.1), In(4.5), In(4.1))
        _add_txt(s, d.get('Adresse',''), In(4.3), In(6.3), In(4.5), In(0.4),
                 size=10, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    else:
        _add_rect(s, In(4.3), In(2.1), In(4.5), In(4.1), GRAY_LIGHT)
        _add_txt(s, "[Adresse + Carte]", In(4.3), In(3.9), In(4.5), In(0.5),
                 size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    # Col 3 : photo bien
    photo = _get_photo(d.get('Photo bien (URL)','') or d.get('photo_url',''))
    if photo:
        _add_img(s, photo, In(9.0), In(2.1), In(4.1), In(4.1))
    else:
        _add_rect(s, In(9.0), In(2.1), In(4.1), In(4.1), GRAY_LIGHT)
        _add_txt(s, "[Photo]", In(9.0), In(3.9), In(4.1), In(0.5),
                 size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(0.2), In(7.08), In(1.2), In(0.38))

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — DESCRIPTION DU BIEN
# ══════════════════════════════════════════════════════════════
def slide_description(prs, d, A):
    s = _blank(prs)
    _add_rect(s, 0, 0, SW, SH, WHITE)

    _add_chevron_shape(s, In(4.8), In(0.12), In(0.5), In(0.5), ORANGE)
    _add_txt(s, "Description du bien", In(5.5), In(0.1), In(7.5), In(0.65),
             size=26, color=TEAL)
    _add_chevron_shape(s, In(12.4), In(0.12), In(0.5), In(0.5), TEAL_DARK)

    photo = _get_photo(d.get('Photo bien (URL)','') or d.get('photo_url',''))
    if photo:
        _add_img(s, photo, In(0), In(0), In(5.0), In(3.75))

    annonce = d.get('Version portail','') or d.get('version_portail','') or '[annonce bien]'
    _add_txt(s, annonce[:600], In(5.5), In(1.0), In(7.5), In(2.65),
             size=12, color=TEAL)

    # Bas
    desc_v = d.get('Description ville','') or d.get('description_ville','') or '[Quartier]'
    _add_rect(s, 0, In(3.75), In(4.4), In(3.7), GRAY_LIGHT)
    _add_txt(s, desc_v[:400], In(0.2), In(3.95), In(4.0), In(3.3),
             size=10, color=DARK)

    _add_rect(s, In(4.5), In(3.75), In(4.3), In(3.7), RGBColor(0xE0,0xE8,0xEE))
    _add_txt(s, "[Photo 3]", In(4.5), In(5.4), In(4.3), In(0.5),
             size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    _add_rect(s, In(8.9), In(3.75), In(4.4), In(3.7), GRAY_LIGHT)
    _add_txt(s, "[Photo 4]", In(8.9), In(5.4), In(4.4), In(0.5),
             size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(0.2), In(7.08), In(1.2), In(0.38))

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — URBANISME & CADASTRE
# ══════════════════════════════════════════════════════════════
def slide_cadastre(prs, d, A, cadastre_buf=None):
    s = _blank(prs)
    _add_rect(s, 0, 0, SW, SH, WHITE)

    _add_chevron_shape(s, In(0.3), In(0.12), In(0.5), In(0.5), ORANGE)
    _add_txt(s, "Urbanisme et zone", In(0.95), In(0.1), In(10.0), In(0.65),
             size=26, color=TEAL)
    _add_chevron_shape(s, In(11.0), In(0.12), In(0.5), In(0.5), TEAL_DARK)

    # Infos cadastrales droite
    section    = d.get('Section cadastrale','') or ''
    parcelle   = d.get('N° parcelle','') or d.get('parcelle','') or ''
    surf_parc  = d.get('Surface parcelle','') or ''
    zone_plu   = d.get('Zone PLU','') or ''
    contraintes= d.get('Contraintes notables','') or ''

    _add_rect(s, In(8.5), In(0.9), In(4.83), In(2.8), TEAL)
    cad_lines = []
    if section:   cad_lines.append(f"Section : {section}")
    if parcelle:  cad_lines.append(f"N° parcelle : {parcelle}")
    if surf_parc: cad_lines.append(f"Surface parcelle : {surf_parc} m²")
    _add_txt(s, "\n".join(cad_lines) if cad_lines else "[Informations cadastrales]",
             In(8.8), In(1.2), In(4.2), In(2.3), size=13, color=WHITE)

    plu_lines = []
    if zone_plu:    plu_lines.append(f"Zone PLU : {zone_plu}")
    if contraintes: plu_lines.append(f"\n{contraintes}")
    _add_txt(s, "\n".join(plu_lines) if plu_lines else "[Information PLU]",
             In(0.3), In(3.8), In(4.2), In(3.4), size=12, color=DARK)

    if cadastre_buf:
        cadastre_buf.seek(0)
        _add_img(s, cadastre_buf, In(4.0), In(1.8), In(4.5), In(5.6))
    else:
        _add_rect(s, In(4.0), In(1.8), In(4.5), In(5.6), GRAY_LIGHT)
        _add_txt(s, "[image plan cadastre]", In(4.0), In(4.4), In(4.5), In(0.5),
                 size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(0.2), In(7.08), In(1.2), In(0.38))

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — VALEUR LOCATIVE
# ══════════════════════════════════════════════════════════════
def slide_valeur(prs, d, A):
    s = _blank(prs)
    _add_rect(s, 0, 0, SW, SH, WHITE)

    _add_chevron_shape(s, In(4.8), In(0.12), In(0.5), In(0.5), ORANGE)
    _add_txt(s, "Valeur locative", In(5.5), In(0.1), In(7.5), In(0.65),
             size=26, color=TEAL)
    _add_chevron_shape(s, In(12.4), In(0.12), In(0.5), In(0.5), TEAL_DARK)

    if A.get('AVIS_VALEUR_B64'):
        _add_img(s, A['AVIS_VALEUR_B64'], In(0), In(0.9), In(5.0), In(2.8))

    avis = d.get('Avis de valeur','') or d.get('avis_de_valeur','') or '[avis de valeur]'
    if '---SYNTH' in avis.upper():
        parts = avis.upper().split('---SYNTH')
        avis_court = '---SYNTH' + parts[1] if len(parts) > 1 else avis
        avis_court = avis_court.split('---')[0].replace('---SYNTHÈSE---','').replace('---SYNTH\u00c8SE---','').strip()
        avis_court = avis.split('---')[1].strip() if '---' in avis else avis[:400]
    else:
        avis_court = avis[:400]

    _add_txt(s, avis_court, In(5.5), In(1.2), In(7.5), In(2.3),
             size=12, color=TEAL)

    _add_rect(s, 0, In(3.75), SW, In(3.7), TEAL)

    if A.get('TAXE_B64'):
        _add_img(s, A['TAXE_B64'], In(4.0), In(3.85), In(5.0), In(3.4))

    taxe = d.get('Taxe foncière','') or d.get('taxe_fonciere','') or ''
    if taxe:
        _add_txt(s, f"Taxe foncière :\n{taxe}", In(9.5), In(4.5), In(3.5), In(2.5),
                 size=16, bold=True, color=WHITE)
    else:
        _add_txt(s, "[Montant de la\nTaxe foncière]", In(9.5), In(5.2), In(3.5), In(1.5),
                 size=14, color=WHITE)

    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(0.2), In(7.08), In(1.2), In(0.38))

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — GALERIE
# ══════════════════════════════════════════════════════════════
def slide_galerie(prs, d, A):
    s = _blank(prs)
    _add_rect(s, 0, 0, SW, SH, TEAL)

    photo = _get_photo(d.get('Photo bien (URL)','') or d.get('photo_url',''))
    if photo:
        _add_img(s, photo, In(0), In(0), SW, In(4.05))

    for i in range(3):
        x = In(0.05) + i * In(4.43)
        _add_rect(s, x, In(4.1), In(4.3), In(3.35), TEAL_DARK)
        label = ["Photo 2", "Photo 3", "Photo 4"][i]
        _add_txt(s, f"[{label}]", x, In(5.5), In(4.3), In(0.5),
                 size=14, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION
# ══════════════════════════════════════════════════════════════
def slide_conclusion(prs, d, A):
    s = _blank(prs)
    _add_rect(s, 0, 0, SW, SH, WHITE)

    if A.get('CONCLUSION_B64'):
        _add_img(s, A['CONCLUSION_B64'], In(0), In(0), In(4.5), SH)

    _add_chevron_shape(s, In(5.0), In(0.12), In(0.5), In(0.5), ORANGE)
    _add_txt(s, "Conclusion", In(5.7), In(0.1), In(7.5), In(0.65),
             size=28, color=TEAL)
    _add_chevron_shape(s, In(12.4), In(0.12), In(0.5), In(0.5), TEAL_DARK)

    prix = d.get('Prix de vente','') or d.get('prix_vente','') or ''
    loyer_m = d.get('Loyer mensuel','') or d.get('loyer_mensuel','') or ''

    _add_txt(s,
        "La valorisation de l'actif pour cet ensemble immobilier à usage commercial est de :",
        In(5.0), In(2.0), In(7.8), In(1.2),
        size=16, bold=True, color=DARK)

    val_display = _pfmt(prix) if prix else (f"{int(float(loyer_m)):,} €/mois" if loyer_m else '[Valeur]')
    _add_txt(s, val_display, In(6.0), In(3.5), In(6.0), In(1.0),
             size=30, color=TEAL, align=PP_ALIGN.CENTER)

    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(0.2), In(7.08), In(1.2), In(0.38))

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — CONTACT
# ══════════════════════════════════════════════════════════════
def slide_contact(prs, d, A):
    s = _blank(prs)
    _add_rect(s, 0, 0, SW, SH, WHITE)
    _add_rect(s, 0, 0, In(4.0), SH, TEAL)
    _add_rect(s, In(9.0), 0, In(4.33), SH/2, TEAL)
    _add_rect(s, In(9.0), SH/2, In(4.33), SH/2, TEAL_DARK)

    if A.get('CONTACT_VOILIER_B64'):
        _add_img(s, A['CONTACT_VOILIER_B64'], In(0), In(0), In(4.0), SH)
    if A.get('CONTACT_VANNES_B64'):
        _add_img(s, A['CONTACT_VANNES_B64'], In(9.0), In(0), In(4.33), SH/2)

    for i, m in enumerate(["Estimation", "Vente", "Location", "Cession d'entreprise"]):
        y = In(0.5) + i * In(0.62)
        indent = i * In(0.22)
        _add_chevron_shape(s, In(0.3)+indent, y+In(0.1), In(0.3), In(0.3), ORANGE)
        _add_txt(s, m, In(0.75)+indent, y, In(3.0), In(0.5),
                 size=14, bold=True, color=WHITE)

    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(4.6), In(1.9), In(3.5), In(1.6))
    _add_txt(s, "contact@barbierimmobilier.com",
             In(4.3), In(4.1), In(4.5), In(0.5),
             size=13, color=ORANGE, align=PP_ALIGN.CENTER)
    _add_txt(s, "barbierimmobilier.com",
             In(4.3), In(4.8), In(4.5), In(0.5),
             size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    nego_tel = d.get('nego_tel','') or '02.97.47.11.11'
    _add_txt(s, "2 place Albert Einstein\n56000 VANNES",
             In(9.2), In(4.0), In(3.9), In(1.0),
             size=16, bold=True, color=WHITE)
    _add_txt(s, f"{nego_tel}\ncontact@barbierimmobilier.com\nbarbierimmobilier.com",
             In(9.2), In(5.2), In(3.9), In(1.6), size=12, color=WHITE)

    if A.get('LOGO_B64'):
        _add_img(s, A['LOGO_B64'], In(0.2), In(7.08), In(1.2), In(0.38))

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING
# ══════════════════════════════════════════════════════════════
def slide_closing(prs, d, A):
    s = _blank(prs)
    _add_rect(s, 0, 0, SW, SH, TEAL)
    if A.get('CLOSING_BG_B64'):
        _add_img(s, A['CLOSING_BG_B64'], In(0), In(0), SW, SH)
    _add_chevron_shape(s, In(1.3), In(1.4), In(0.5), In(0.5), ORANGE)
    _add_txt(s, "Votre projet", In(2.0), In(1.3), In(5.0), In(1.0),
             size=32, color=TEAL, name="Calibri Light")
    _add_txt(s, "devient le nôtre", In(6.0), In(4.9), In(6.0), In(1.0),
             size=32, color=TEAL, name="Calibri Light")
    _add_chevron_shape(s, In(5.4), In(5.0), In(0.55), In(0.55), ORANGE)

# ══════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════
def generate_dossier_pptx(d, assets, map_buf=None, cadastre_buf=None):
    prs = _prs()
    A = assets
    slide_couverture(prs, d, A)
    slide_sommaire(prs, d, A)
    slide_agence(prs, d, A)
    slide_equipe(prs, d, A)
    slide_secteur(prs, d, A, map_buf=map_buf)
    slide_description(prs, d, A)
    slide_cadastre(prs, d, A, cadastre_buf=cadastre_buf)
    slide_valeur(prs, d, A)
    slide_galerie(prs, d, A)
    slide_conclusion(prs, d, A)
    slide_contact(prs, d, A)
    slide_closing(prs, d, A)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

print("gen_pptx_v2.py OK")
